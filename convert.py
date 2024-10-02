import os
import sys
import pypandoc
import io
from openpyxl import load_workbook
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, landscape
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Image, Spacer
from reportlab.lib.units import inch
from reportlab.platypus import PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from pptx import Presentation
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPDF
import pdfkit

def download_pdf():
    return os.path.join(os.path.expanduser('~'), 'Downloads')

def generate_unique_filename(base_path, base_name, extension):
    counter = 1
    new_path = os.path.join(base_path, f"{base_name}{extension}")
    while os.path.exists(new_path):
        new_path = os.path.join(base_path, f"{base_name}_{counter}{extension}")
        counter += 1
    return new_path

#consversor de DOCX
def convert_docx_to_pdf(input_path: str, output_path: str) -> None:
    output_path = output_path if output_path.endswith('.pdf') else output_path + '.pdf'
    try:
        pypandoc.convert_file(input_path, 'pdf', outputfile=output_path, 
                              extra_args=['--pdf-engine=/Library/TeX/texbin/pdflatex'])
    except Exception as e:
        print(f"Error converting with pdflatex: {e}")
        print("Attempting conversion via HTML...")
        try:
            html_content = pypandoc.convert_file(input_path, 'html')
            pdfkit.from_string(html_content, output_path)
        except Exception as e:
            print(f"Error converting via HTML: {e}")
            raise

#conversor XLSX
# Función para obtener anchos de columnas
def get_column_widths(sheet, max_width):
    column_widths = []
    for col in sheet.columns:
        width = sheet.column_dimensions[col[0].column_letter].width
        if width is None:
            width = 10
        column_widths.append(width * 8)
    
    total_width = sum(column_widths)
    if total_width > max_width:
        scale = max_width / total_width
        column_widths = [width * scale for width in column_widths]
    
    return column_widths

# Función para obtener la altura de las filas y limitarla si es necesario
def get_row_heights(sheet, max_height=50):
    row_heights = []
    for row in sheet.iter_rows():
        height = sheet.row_dimensions[row[0].row].height
        if height is None:
            height = 15
        if height > max_height:  # Limitar altura
            height = max_height
        row_heights.append(height)
    return row_heights

# Función para procesar celdas y ajustarlas a tamaños razonables
def process_cell(cell, centered_style, max_characters=100):
    if cell.data_type == 'i':  # Image
        img = Image(io.BytesIO(cell.value))
        img.drawHeight = 50  # Ajusta según sea necesario
        img.drawWidth = 50   # Ajusta según sea necesario
        return img
    elif cell.value is not None:
        cell_value = str(cell.value)
        if len(cell_value) > max_characters:  # Limitar longitud del contenido
            cell_value = cell_value[:max_characters] + '...'  # Truncar contenido largo
        return Paragraph(cell_value, centered_style)
    return ''

# Función para convertir XLSX a PDF
def convert_xlsx_to_pdf(input_path: str, output_path: str) -> None:
    wb = load_workbook(input_path)
    sheet = wb.active
    
    # Configurar márgenes y tamaño de página
    page_width, page_height = landscape(letter)
    margins = 0.5 * inch
    effective_page_width = page_width - 2 * margins
    effective_page_height = page_height - 2 * margins
    
    pdf = SimpleDocTemplate(
        output_path, pagesize=landscape(letter), 
        leftMargin=margins, rightMargin=margins, topMargin=margins, bottomMargin=margins
    )
    
    data = []
    styles = getSampleStyleSheet()
    centered_style = ParagraphStyle('centered', parent=styles['Normal'], alignment=TA_CENTER, wordWrap='CJK')
    
    for row in sheet.iter_rows():
        processed_row = []
        for cell in row:
            processed_cell = process_cell(cell, centered_style)
            processed_row.append(processed_cell)
        data.append(processed_row)
    
    # Ajustar anchos de las columnas
    column_widths = get_column_widths(sheet, effective_page_width)
    
    # Definir la tabla con los datos procesados
    table = Table(data, colWidths=column_widths, repeatRows=1)
    
    # Estilos de la tabla
    style = TableStyle([
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('LEFTPADDING', (0, 0), (-1, -1), 3),
        ('RIGHTPADDING', (0, 0), (-1, -1), 3),
        ('TOPPADDING', (0, 0), (-1, -1), 3),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 3),
        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
        ('FONTSIZE', (0, 0), (-1, -1), 8),  # Ajuste de tamaño de fuente
    ])
    
    # Combinar celdas
    for merged_range in sheet.merged_cells.ranges:
        min_col, min_row, max_col, max_row = merged_range.bounds
        style.add('SPAN', (min_col - 1, min_row - 1), (max_col - 1, max_row - 1))
    
    table.setStyle(style)
    elements = [table]
    
    # Intentar construir el PDF
    try:
        pdf.build(elements)
        #print(f'Archivo PDF guardado en: {output_path}')
    except Exception as e:
        raise ValueError(f"Error during conversion: {e}")

#conversor PPT
# def convert_pptx_to_pdf(input_path: str, output_path: str) -> None:
    prs = Presentation(input_path)
    slides_text = "\n\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    pypandoc.convert_text(slides_text, 'pdf', format='markdown', outputfile=output_path)
def extract_text_from_shape(shape):
    """Extrae el texto de un shape y respeta el formato (negritas, listas)."""
    text = ""
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.bold:
                    text += f"<b>{run.text}</b>"
                else:
                    text += run.text
            text += "\n"
    return text

def extract_image_from_shape(shape, slide_number, output_folder):
    """Extrae la imagen de un shape si es un placeholder o imagen."""
    if shape.shape_type == 13:  # Es una imagen
        image = shape.image
        image_filename = f"slide_{slide_number}_image_{shape.shape_id}.png"
        image_path = os.path.join(output_folder, image_filename)
        with open(image_path, "wb") as img_file:
            img_file.write(image.blob)
        return image_path
    return None

def extract_table_from_shape(shape):
    """Extrae tablas de los shapes si hay alguna en la diapositiva."""
    if shape.has_table:
        table_data = []
        for row in shape.table.rows:
            row_data = [cell.text for cell in row.cells]
            table_data.append(row_data)
        return table_data
    return None

def convert_pptx_to_pdf(input_path: str, output_path: str):
    prs = Presentation(input_path)
    pdf = SimpleDocTemplate(output_path, pagesize=landscape(letter))

    elements = []
    styles = getSampleStyleSheet()
    style_normal = styles['Normal']
    style_heading = styles['Heading1']

    output_images_folder = os.path.join(os.path.dirname(output_path), "pptx_images")
    os.makedirs(output_images_folder, exist_ok=True)

    for slide_idx, slide in enumerate(prs.slides):
        elements.append(Paragraph(f"Slide {slide_idx + 1}", style_heading))
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # Extraer el texto formateado
                text = extract_text_from_shape(shape)
                if text:
                    elements.append(Paragraph(text, style_normal))
            
            # Extraer las imágenes
            image_path = extract_image_from_shape(shape, slide_idx + 1, output_images_folder)
            if image_path:
                img = Image(image_path, width=4*inch, height=3*inch)
                elements.append(img)

            # Extraer tablas
            table_data = extract_table_from_shape(shape)
            if table_data:
                table = Table(table_data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ]))
                elements.append(table)

        elements.append(Spacer(1, 12))  # Añadir espacio entre slides
        elements.append(PageBreak())  # Nueva página para cada diapositiva

    pdf.build(elements)
    print(f"PDF creado exitosamente en {output_path}")

def convert_txt_to_pdf(input_path: str, output_path: str) -> None:
    with open(input_path, 'r', encoding='utf-8') as f:
        content = f.read()
    pypandoc.convert_text(content, 'pdf', format='markdown', outputfile=output_path)

def convert_svg_to_pdf(input_path: str, output_path: str) -> None:
    drawing = svg2rlg(input_path)
    renderPDF.drawToFile(drawing, output_path)

def convert_file_to_pdf(input_path: str, output_path: str) -> None:
    try:
        ext = os.path.splitext(input_path)[1].lower()
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        
        downloads = download_pdf()
        output_path = generate_unique_filename(downloads, base_name, '.pdf')
        
        if ext == '.docx':
            convert_docx_to_pdf(input_path, output_path)
        elif ext == '.xlsx':
            convert_xlsx_to_pdf(input_path, output_path)
        elif ext == '.pptx':
            convert_pptx_to_pdf(input_path, output_path)
        elif ext == '.txt':
            convert_txt_to_pdf(input_path, output_path)
        elif ext == '.svg':
            convert_svg_to_pdf(input_path, output_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
        
        print(f"{output_path}")
        return output_path
    except Exception as e:
        error_message = f"Error during conversion: {str(e)}"
        print(error_message)
        raise ValueError(error_message)

# if __name__ == '__main__':
#     input_file = '/Users/dev13/desktop/Buscador de Dioses.docx'  
#     output_file = 'sisas.pdf' 
#     convert_file_to_pdf(input_file, output_file)
    
if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Uso: python convert.py <input_file>")
        sys.exit(1)
    
    input_file = sys.argv[1]
    convert_file_to_pdf(input_file, "")
    