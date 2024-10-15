from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.colors import HexColor
from io import BytesIO
from PIL import Image as PILImage
from formats.pptx_fonts import extract_text_properties, get_shape_position
from reportlab.lib.styles import ParagraphStyle
from reportlab.platypus import Paragraph
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
from reportlab.lib.colors import black
from reportlab.lib.units import inch

def extract_image_from_shape(shape, pdf, slide_height):
    if shape.shape_type == 13: 
        img_stream = BytesIO(shape.image.blob)
        img = PILImage.open(img_stream)
        left, top, width, height = get_shape_position(shape, slide_height)
        pdf.drawInlineImage(img, left, top - height, width, height)

def extract_background_image(slide):
    background = slide.background
    if background.fill.type == 6:  # Background con imagen
        image = background.fill.picture.image
        image_stream = BytesIO(image.blob)
        return PILImage.open(image_stream)
    return None

def extract_table_from_shape(shape):
    if shape.has_table:
        table_data = []
        for row in shape.table.rows:
            row_data = [cell.text for cell in row.cells]
            table_data.append(row_data)
        return table_data
    return None

def draw_text_with_properties(pdf, text_props, left, top, width, height):
    total_height = sum(prop.get('font_size', 12) * 1.2 for prop in text_props)
    
    # Start from the top of the text box
    current_y = top
    
    vertical_alignment = text_props[0].get('vertical_alignment', 0)
    if vertical_alignment == 1:  # Middle
        current_y = top - (height - total_height) / 2
    elif vertical_alignment == 2:  # Bottom
        current_y = top - height + total_height
    
    for prop in text_props:
        font_name = prop.get('font_name', 'Helvetica')
        if prop.get('bold') and prop.get('italic'):
            font_name += '-BoldOblique'
        elif prop.get('bold'):
            font_name += '-Bold'
        elif prop.get('italic'):
            font_name += '-Oblique'

        alignment = TA_LEFT
        if prop.get('alignment') == 1:
            alignment = TA_CENTER
        elif prop.get('alignment') == 2:
            alignment = TA_RIGHT
        elif prop.get('alignment') == 3:
            alignment = TA_JUSTIFY

        text_color = black
        if 'color' in prop:
            try:
                text_color = HexColor(prop['color'])
            except ValueError:
                print(f"Warning: Invalid color value '{prop['color']}'. Using black instead.")

        font_size = prop.get('font_size', 12)
        style = ParagraphStyle(
            name='Custom',
            fontName=font_name,
            fontSize=font_size,
            textColor=text_color,
            alignment=alignment,
            leading=font_size * 1.2
        )

        para = Paragraph(prop.get('text', ''), style)
        para.wrapOn(pdf, width, height)
        para.drawOn(pdf, left, current_y - para.height)
        current_y -= para.height

    return current_y

def convert_pptx_to_pdf(input_path: str, output_path: str):
    prs = Presentation(input_path)
    slide_width = prs.slide_width / 914400 * inch
    slide_height = prs.slide_height / 914400 * inch
    
    pdf = canvas.Canvas(output_path, pagesize=(slide_width, slide_height))
    
    for slide in prs.slides:
        fill = slide.background.fill
        if fill.type == 1:  # Solid fill
            bg_color = f"#{fill.fore_color.rgb:06x}"
            pdf.setFillColor(HexColor(bg_color))
            pdf.rect(0, 0, slide_width, slide_height, fill=1)

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_props = extract_text_properties(shape, slide.element.xml, slide.part.slide.element.nsmap)
                if text_props:
                    left, top, width, height = get_shape_position(shape, slide_height)
                    draw_text_with_properties(pdf, text_props, left, top, width, height)
            elif shape.shape_type == 13:  # Picture
                extract_image_from_shape(shape, pdf, slide_height)
        
        pdf.showPage()
    
    pdf.save()